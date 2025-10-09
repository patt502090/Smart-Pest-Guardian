from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def build_presentation(output_path: str) -> None:
    prs = Presentation()

    # Slide 1: Table summary
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "การต่อยอดบทเรียนสัปดาห์ 9–14 เข้าสู่โปรเจกต์"

    rows, cols = 7, 3
    table_width, table_height = Inches(10), Inches(3.5)
    left, top = Inches(0.5), Inches(1.5)
    table = slide.shapes.add_table(rows, cols, left, top, table_width, table_height).table

    headers = ["Week", "หัวข้อในคอร์ส", "สิ่งที่หยิบมาใช้ในโปรเจกต์"]
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(20)

    entries = [
        ("Week 9", "Neural Networks", "Backbone ลึกหลายชั้นของ YOLOv8 + การใช้ box/cls/dfl loss"),
        ("Week 10", "CNNs", "โครงสร้าง CNN-based detector (CSP, Focus, SPPF) สำหรับดึงฟีเจอร์"),
        ("Week 11", "RNNs", "บันทึกไว้ในรายงานว่าไม่ใช้ RNN เพราะโจทย์นี้เน้นภาพนิ่ง"),
        ("Week 12", "NLP & Transformers", "เทียบให้เห็นว่าเราเลือก Vision เป็นหลัก แต่ชี้แนวทางต่อยอดด้วย Transformer"),
        ("Week 13", "Reinforcement Learning", "เสนอ future work ให้ RL ช่วยปรับ threshold จาก feedback เกษตรกร"),
        ("Week 14", "Model Evaluation & Optimization", "ใช้ mAP, Precision/Recall, F1 curve และ early stopping ลด overfitting"),
    ]

    for row_idx, (week, topic, usage) in enumerate(entries, start=1):
        table.cell(row_idx, 0).text = week
        table.cell(row_idx, 1).text = topic
        table.cell(row_idx, 2).text = usage
        for col in range(cols):
            para = table.cell(row_idx, col).text_frame.paragraphs[0]
            para.font.size = Pt(16)

    # Slide 2: Timeline
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Timeline: นำบทเรียนสู่ Smart Pest Guardian"

    start_left = Inches(0.7)
    start_top = Inches(1.8)
    step = Inches(1.6)
    circle_size = Inches(0.7)
    line_height = Inches(0.05)

    weeks = [
        ("W9", "Neural\nNetworks", "Backbone + Loss"),
        ("W10", "CNNs", "Feature extractor"),
        ("W11", "RNNs", "Note vision focus"),
        ("W12", "NLP/\nTransformers", "Future hybrid"),
        ("W13", "RL", "Auto threshold"),
        ("W14", "Eval & Opt", "mAP / tuning"),
    ]

    for idx, (short, title, note) in enumerate(weeks):
        left = start_left + step * idx
        top = start_top

        if idx > 0:
            line_left = start_left + step * idx - step / 2
            connector = slide2.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                line_left,
                top + circle_size / 2 - line_height / 2,
                step,
                line_height,
            )
            connector.fill.solid()
            connector.fill.fore_color.rgb = RGBColor(0, 120, 212)
            connector.line.fill.background()

        circle = slide2.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            left,
            top,
            circle_size,
            circle_size,
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(0, 120, 212)
        circle.line.color.rgb = RGBColor(255, 255, 255)
        circle.line.width = Pt(2)

        tf = circle.text_frame
        tf.text = short
        para = tf.paragraphs[0]
        para.font.size = Pt(20)
        para.font.bold = True
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER

        textbox = slide2.shapes.add_textbox(
            left - Inches(0.15),
            top + circle_size + Inches(0.1),
            circle_size + Inches(0.3),
            Inches(1.1),
        )
        tf_box = textbox.text_frame
        tf_box.word_wrap = True
        p_title = tf_box.paragraphs[0]
        p_title.text = title
        p_title.font.bold = True
        p_title.font.size = Pt(16)
        p_title.font.color.rgb = RGBColor(0, 51, 102)
        p_title.alignment = PP_ALIGN.CENTER

        p_note = tf_box.add_paragraph()
        p_note.text = note
        p_note.font.size = Pt(14)
        p_note.font.color.rgb = RGBColor(51, 51, 51)
        p_note.alignment = PP_ALIGN.CENTER

    caption = slide2.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.6))
    caption_tf = caption.text_frame
    caption_tf.text = "การต่อยอดบทเรียนสัปดาห์ 9–14 เข้าสู่โปรเจกต์ Smart Pest Guardian"
    caption_para = caption_tf.paragraphs[0]
    caption_para.font.size = Pt(16)
    caption_para.font.color.rgb = RGBColor(80, 80, 80)
    caption_para.alignment = PP_ALIGN.CENTER

    prs.save(output_path)


if __name__ == "__main__":
    build_presentation("reports/section3_week_linkage.pptx")
